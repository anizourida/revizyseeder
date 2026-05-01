#!/usr/bin/env python3
import argparse
import collections
import collections.abc
import hashlib
import json
import os
import shutil
import sys
import zipfile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def ensure_dir(directory: str) -> None:
    os.makedirs(directory, exist_ok=True)


def calculate_blob_hash(blob: bytes) -> str:
    return hashlib.md5(blob).hexdigest()


def convert_ppsx_to_pptx(ppsx_path: str, pptx_path: str) -> bool:
    try:
        with zipfile.ZipFile(ppsx_path, "r") as zin:
            with zipfile.ZipFile(pptx_path, "w") as zout:
                for item in zin.infolist():
                    buffer = zin.read(item.filename)
                    if item.filename == "[Content_Types].xml":
                        content = buffer.decode("utf-8")
                        content = content.replace(
                            "application/vnd.openxmlformats-officedocument.presentationml.slideshow.main+xml",
                            "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml",
                        )
                        zout.writestr(item, content.encode("utf-8"))
                    else:
                        zout.writestr(item, buffer)
        return True
    except Exception:
        return False


def extract_image(shape, slide_id: int, image_count: int, assets_dir: str, dedupe_map: dict) -> str | None:
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_bytes = shape.image.blob
            image_ext = shape.image.ext
        elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and hasattr(shape, "image"):
            image_bytes = shape.image.blob
            image_ext = shape.image.ext
        else:
            return None

        file_hash = calculate_blob_hash(image_bytes)
        if file_hash in dedupe_map:
            return dedupe_map[file_hash]

        image_name = f"slide_{slide_id}_image_{image_count}.{image_ext}"
        image_path = os.path.join(assets_dir, image_name)
        with open(image_path, "wb") as f:
            f.write(image_bytes)

        dedupe_map[file_hash] = image_name
        return image_name
    except Exception:
        return None


def extract_video_from_zip(shape, slide, pptx_path: str, slide_id: int, video_count: int, assets_dir: str) -> str | None:
    try:
        ns = {
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        }

        nv_pic_pr = shape._element.find(".//p:nvPicPr", ns)
        if nv_pic_pr is None:
            return None

        nv_pr = nv_pic_pr.find(".//p:nvPr", ns)
        if nv_pr is None:
            return None

        video_file = nv_pr.find(".//a:videoFile", ns)
        if video_file is None:
            return None

        rid = video_file.get(f"{{{ns['r']}}}link")
        if not rid or rid not in slide.part.rels:
            return None

        rel = slide.part.rels[rid]
        if rel.is_external:
            return None

        internal_path = rel.target_part.partname
        if internal_path.startswith("/"):
            internal_path = internal_path[1:]

        ext = os.path.splitext(internal_path)[1]
        video_filename = f"slide_{slide_id}_video_{video_count}{ext}"
        output_path = os.path.join(assets_dir, video_filename)

        with zipfile.ZipFile(pptx_path, "r") as z:
            with z.open(internal_path) as source, open(output_path, "wb") as target:
                shutil.copyfileobj(source, target)

        return video_filename
    except Exception:
        return None


def process_presentation(file_path: str, lesson_id: str, output_root: str) -> dict:
    lesson_dir = os.path.join(output_root, lesson_id)
    assets_dir = os.path.join(lesson_dir, "assets")
    ensure_dir(lesson_dir)

    if os.path.exists(assets_dir):
        shutil.rmtree(assets_dir)
    ensure_dir(assets_dir)

    temp_pptx = os.path.join(lesson_dir, "temp_processing.pptx")

    is_ppsx = file_path.lower().endswith(".ppsx")
    if is_ppsx:
        if not convert_ppsx_to_pptx(file_path, temp_pptx):
            raise RuntimeError("Failed to convert PPSX to PPTX")
    else:
        shutil.copy2(file_path, temp_pptx)

    try:
        prs = Presentation(temp_pptx)
    except Exception as e:
        if os.path.exists(temp_pptx):
            os.remove(temp_pptx)
        raise RuntimeError(f"Could not open presentation: {e}") from e

    lesson_data = {
        "file_name": os.path.basename(file_path),
        "lesson_id": lesson_id,
        "metadata": {
            "total_slides": len(prs.slides),
            "slide_width_emu": int(prs.slide_width),
            "slide_height_emu": int(prs.slide_height),
        },
        "slides": [],
    }

    dedupe_map = {}

    for i, slide in enumerate(prs.slides):
        slide_id = i + 1
        slide_data = {"id": slide_id, "elements": []}
        image_count = 0
        video_count = 0

        def extract_recursive(shapes, parent_offset=(0, 0), parent_group_context=None):
            nonlocal image_count
            nonlocal video_count

            for shape in list(shapes):
                try:
                    shape_type = shape.shape_type
                    abs_left = int(shape.left) + int(parent_offset[0])
                    abs_top = int(shape.top) + int(parent_offset[1])
                    bbox = [abs_left, abs_top, int(shape.width), int(shape.height)]
                except Exception:
                    # Some malformed/unsupported shapes crash python-pptx when accessed.
                    # Skip these shapes and continue extracting other slide elements.
                    continue

                try:
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            element = {
                                "type": "text",
                                "content": text,
                                "bbox": bbox,
                            }
                            if parent_group_context:
                                element["group_context"] = parent_group_context
                            slide_data["elements"].append(element)
                except Exception:
                    pass

                if shape_type in [MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.PLACEHOLDER]:
                    image_count += 1
                    image_filename = extract_image(
                        shape=shape,
                        slide_id=slide_id,
                        image_count=image_count,
                        assets_dir=assets_dir,
                        dedupe_map=dedupe_map,
                    )
                    if image_filename:
                        slide_data["elements"].append(
                            {
                                "type": "image",
                                "file_path": os.path.join("assets", image_filename),
                                "bbox": bbox,
                                "description": shape.name,
                            }
                        )

                if shape_type == MSO_SHAPE_TYPE.MEDIA:
                    video_count += 1
                    video_filename = extract_video_from_zip(
                        shape=shape,
                        slide=slide,
                        pptx_path=temp_pptx,
                        slide_id=slide_id,
                        video_count=video_count,
                        assets_dir=assets_dir,
                    )
                    if video_filename:
                        slide_data["elements"].append(
                            {
                                "type": "video",
                                "file_path": os.path.join("assets", video_filename),
                                "bbox": bbox,
                                "description": shape.name,
                            }
                        )

                if shape_type == MSO_SHAPE_TYPE.GROUP:
                    try:
                        extract_recursive(
                            shapes=shape.shapes,
                            parent_offset=(abs_left, abs_top),
                            parent_group_context=shape.name or parent_group_context,
                        )
                    except Exception:
                        pass

        extract_recursive(slide.shapes)
        lesson_data["slides"].append(slide_data)

    json_path = os.path.join(lesson_dir, "data.json")
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(lesson_data, f, indent=2, ensure_ascii=False)

    if os.path.exists(temp_pptx):
        os.remove(temp_pptx)

    image_count = 0
    video_count = 0
    for slide in lesson_data["slides"]:
        for element in slide["elements"]:
            if element.get("type") == "image":
                image_count += 1
            elif element.get("type") == "video":
                video_count += 1

    return {
        "lesson_id": lesson_id,
        "json_path": json_path,
        "assets_dir": assets_dir,
        "total_slides": int(lesson_data["metadata"]["total_slides"]),
        "images": image_count,
        "videos": video_count,
    }


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Extract slide JSON + assets from PPTX/PPSX")
    parser.add_argument("--input", required=True, help="Absolute path to PPTX/PPSX file")
    parser.add_argument("--lesson-id", required=False, help="Lesson id for output directory naming")
    parser.add_argument("--output-root", required=True, help="Root output directory")
    return parser.parse_args()


def main() -> int:
    args = parse_args()

    input_path = os.path.abspath(args.input)
    if not os.path.isfile(input_path):
        sys.stderr.write(f"Input file not found: {input_path}\n")
        return 2

    lesson_id = args.lesson_id or os.path.splitext(os.path.basename(input_path))[0]
    output_root = os.path.abspath(args.output_root)
    ensure_dir(output_root)

    try:
        summary = process_presentation(input_path, lesson_id, output_root)
    except Exception as e:
        sys.stderr.write(f"Extraction failed: {e}\n")
        return 1

    sys.stdout.write(json.dumps(summary, ensure_ascii=False) + "\n")
    return 0


if __name__ == "__main__":
    sys.exit(main())
